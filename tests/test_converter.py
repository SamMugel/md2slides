"""Tests for the PPTX converter."""

import os
import tempfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.oxml.ns import qn

from md2slides.converter import MarkdownToPptxConverter, convert_file
from md2slides.parser import ValidationError


class TestConverterValidation:
    """Test converter input validation."""

    def test_empty_output_path_raises_error(self):
        """Empty output path should raise ValidationError."""
        converter = MarkdownToPptxConverter()
        with pytest.raises(ValidationError, match="cannot be empty"):
            converter.convert("# Title", "")

    def test_whitespace_output_path_raises_error(self):
        """Whitespace-only output path should raise ValidationError."""
        converter = MarkdownToPptxConverter()
        with pytest.raises(ValidationError, match="cannot be empty"):
            converter.convert("# Title", "   ")

    def test_non_string_output_path_raises_error(self):
        """Non-string output path should raise ValidationError."""
        converter = MarkdownToPptxConverter()
        with pytest.raises(ValidationError, match="must be a string"):
            converter.convert("# Title", 123)  # type: ignore

    def test_wrong_extension_raises_error(self):
        """Non-.pptx extension should raise ValidationError."""
        converter = MarkdownToPptxConverter()
        with pytest.raises(ValidationError, match=".pptx extension"):
            converter.convert("# Title", "output.pdf")


class TestConverterOutput:
    """Test converter PPTX output."""

    def test_creates_pptx_file(self):
        """Converter should create a valid PPTX file."""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            result = converter.convert("# Test Title", output_path)

            assert os.path.exists(result)
            # Verify it's a valid PPTX
            prs = Presentation(result)
            assert len(prs.slides) == 1

    def test_returns_absolute_path(self):
        """Converter should return absolute path."""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            result = converter.convert("# Test", output_path)

            assert os.path.isabs(result)

    def test_creates_output_directory(self):
        """Converter should create output directory if needed."""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "subdir", "nested", "test.pptx")
            result = converter.convert("# Test", output_path)

            assert os.path.exists(result)


class TestSlideContent:
    """Test slide content generation."""

    def test_title_slide_content(self):
        """Title slide should have correct title."""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert("# My Presentation", output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            # Find text in shapes
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)

            assert "My Presentation" in texts

    def test_title_slide_with_subtitle(self):
        """Title slide should include subtitle."""
        content = """# Main Title

This is the subtitle
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)

            assert "Main Title" in texts
            assert "This is the subtitle" in texts

    def test_content_slide_title(self):
        """Content slide should have correct title."""
        content = """# Title

## Section Header

Some content
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            assert len(prs.slides) == 2

            # Check second slide has the section header
            slide = prs.slides[1]
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)

            assert any("Section Header" in t for t in texts)

    def test_bullet_list_content(self):
        """Bullet list should appear in content slide."""
        content = """## Slide

- First item
- Second item
- Third item
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            # Find content shape
            all_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text += shape.text

            assert "First item" in all_text
            assert "Second item" in all_text
            assert "Third item" in all_text

    def test_multiple_slides(self):
        """Multiple H2s should create multiple slides."""
        content = """# Title

## Slide 1

Content 1

## Slide 2

Content 2

## Slide 3

Content 3
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            assert len(prs.slides) == 4  # 1 title + 3 content


class TestConvertFile:
    """Test the convert_file function."""

    def test_convert_file_basic(self):
        """convert_file should work with basic markdown file."""
        with tempfile.TemporaryDirectory() as tmpdir:
            # Create input file
            input_path = os.path.join(tmpdir, "input.md")
            with open(input_path, "w") as f:
                f.write("# Test Presentation\n\n## Slide One\n\nContent")

            output_path = os.path.join(tmpdir, "output.pptx")
            result = convert_file(input_path, output_path)

            assert os.path.exists(result)
            prs = Presentation(result)
            assert len(prs.slides) == 2

    def test_convert_file_default_output(self):
        """convert_file should use default output path."""
        with tempfile.TemporaryDirectory() as tmpdir:
            input_path = os.path.join(tmpdir, "presentation.md")
            with open(input_path, "w") as f:
                f.write("# Test\n\n## Content\n\nText")

            result = convert_file(input_path)

            expected_output = os.path.join(tmpdir, "presentation.pptx")
            assert result == os.path.abspath(expected_output)
            assert os.path.exists(result)

    def test_convert_file_nonexistent_raises_error(self):
        """convert_file should raise FileNotFoundError for missing file."""
        with pytest.raises(FileNotFoundError, match="not found"):
            convert_file("/nonexistent/path/file.md")

    def test_convert_file_empty_path_raises_error(self):
        """convert_file should raise ValidationError for empty path."""
        with pytest.raises(ValidationError, match="cannot be empty"):
            convert_file("")

    def test_convert_file_non_string_raises_error(self):
        """convert_file should raise ValidationError for non-string path."""
        with pytest.raises(ValidationError, match="must be a string"):
            convert_file(123)  # type: ignore

    def test_convert_file_directory_raises_error(self):
        """convert_file should raise ValidationError for directory path."""
        with tempfile.TemporaryDirectory() as tmpdir:
            with pytest.raises(ValidationError, match="not a file"):
                convert_file(tmpdir)


class TestFormattingPreservation:
    """Test that text formatting is preserved in output."""

    def test_bold_text_preserved(self):
        """Bold formatting should be preserved in PPTX."""
        content = """## Slide

- This has **bold** text
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            # Find content shape and check for bold runs
            has_bold = False
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.bold and "bold" in run.text:
                                has_bold = True

            assert has_bold is True

    def test_italic_text_preserved(self):
        """Italic formatting should be preserved in PPTX."""
        content = """## Slide

- This has *italic* text
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            # Find content shape and check for italic runs
            has_italic = False
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.italic and "italic" in run.text:
                                has_italic = True

            assert has_italic is True


class TestListFormatting:
    """Test that list formatting renders properly in PPTX."""

    def test_bullet_points_have_bullet_char(self):
        """Bullet points should have proper bullet character in XML."""
        content = """## Slide

- First bullet
- Second bullet
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            # Find paragraphs with bullet characters
            bullet_count = 0
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        pPr = para._p.pPr
                        assert pPr is not None
                        buChar = pPr.find(qn('a:buChar'))
                        if buChar is not None:
                            bullet_count += 1

            assert bullet_count == 2

    def test_nested_bullets_have_different_chars(self):
        """Nested bullet points should have different bullet characters."""
        content = """## Slide

- Parent item
  - Child item
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            bullet_chars = []
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        pPr = para._p.pPr
                        assert pPr is not None
                        buChar = pPr.find(qn('a:buChar'))
                        if buChar is not None:
                            bullet_chars.append(buChar.get('char'))

            assert len(bullet_chars) == 2
            # Parent and child should have different bullet chars
            assert bullet_chars[0] != bullet_chars[1]

    def test_numbered_list_has_auto_numbering(self):
        """Numbered lists should use buAutoNum element."""
        content = """## Slide

1. First step
2. Second step
3. Third step
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            auto_num_count = 0
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        pPr = para._p.pPr
                        assert pPr is not None
                        buAutoNum = pPr.find(qn('a:buAutoNum'))
                        if buAutoNum is not None:
                            auto_num_count += 1
                            # Should be arabicPeriod type (1. 2. 3.)
                            assert buAutoNum.get('type') == 'arabicPeriod'

            assert auto_num_count == 3

    def test_nested_numbered_list_uses_alpha(self):
        """Nested numbered lists should use alphabetic numbering."""
        content = """## Slide

1. Parent step
   1. Child step
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            num_types = []
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        pPr = para._p.pPr
                        assert pPr is not None
                        buAutoNum = pPr.find(qn('a:buAutoNum'))
                        if buAutoNum is not None:
                            num_types.append(buAutoNum.get('type'))

            assert len(num_types) == 2
            assert num_types[0] == 'arabicPeriod'  # Parent: 1. 2. 3.
            assert num_types[1] == 'alphaLcPeriod'  # Child: a. b. c.

    def test_mixed_list_formatting(self):
        """Mixed lists should have correct bullet/number formatting."""
        content = """## Slide

- Bullet item
  1. Numbered sub-item
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            has_bullet = False
            has_number = False
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        pPr = para._p.pPr
                        assert pPr is not None
                        buChar = pPr.find(qn('a:buChar'))
                        buAutoNum = pPr.find(qn('a:buAutoNum'))
                        if buChar is not None:
                            has_bullet = True
                        if buAutoNum is not None:
                            has_number = True

            assert has_bullet is True
            assert has_number is True


class TestBrandStyling:
    """Test Multiverse Computing brand styling."""

    def test_slide_has_brand_background_color(self):
        """Slides should have Catskill White background (#F8FAFC)."""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert("# Test Title", output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            # Check background color
            bg_fill = slide.background.fill
            assert bg_fill.fore_color.rgb == (0xF8, 0xFA, 0xFC)

    def test_title_uses_brand_text_color(self):
        """Title text should use Woodsmoke color (#111417)."""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert("# Test Title", output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            # Find title text and check color
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if "Test Title" in run.text:
                                assert run.font.color.rgb == (0x11, 0x14, 0x17)

    def test_title_uses_montserrat_font(self):
        """Title text should use Montserrat font."""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert("# Test Title", output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if "Test Title" in run.text:
                                assert run.font.name == "Montserrat"

    def test_body_uses_open_sans_font(self):
        """Body text should use Open Sans font."""
        content = """## Slide

- Body text here
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if "Body text here" in run.text:
                                assert run.font.name == "Open Sans"

    def test_h1_uses_correct_size(self):
        """H1 title should use 28pt font size (issue #3)."""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert("# Test Title", output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if "Test Title" in run.text:
                                assert run.font.size.pt == 28

    def test_h2_uses_correct_size(self):
        """H2 title should use 28pt font size (issue #3)."""
        content = """## Section Header

Content here
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if "Section Header" in run.text:
                                assert run.font.size.pt == 28

    def test_body_uses_correct_size(self):
        """Body text should use 18pt font size (issue #4)."""
        content = """## Slide

- Body text here
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if "Body text here" in run.text:
                                assert run.font.size.pt == 18


class TestStyleEnhancements:
    """Test style enhancements from issue #3."""

    def test_child_items_use_dark_grey(self):
        """Child list items should use dark grey color (#404040)."""
        content = """## Slide

- Parent item
  - Child item
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            found_child = False
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if "Child item" in run.text:
                                assert run.font.color.rgb == (0x40, 0x40, 0x40)
                                found_child = True
            assert found_child is True

    def test_parent_items_use_woodsmoke(self):
        """Parent list items should use Woodsmoke color (#111417)."""
        content = """## Slide

- Parent item
  - Child item
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            found_parent = False
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if "Parent item" in run.text:
                                assert run.font.color.rgb == (0x11, 0x14, 0x17)
                                found_parent = True
            assert found_parent is True

    def test_list_items_have_no_leading_spaces(self):
        """List items should NOT have leading spaces (issue #1)."""
        content = """## Slide

- Test item
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            # Find the content shape and verify no spacing-only runs
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        runs = list(para.runs)
                        for run in runs:
                            # No run should be only spaces
                            assert run.text.strip() != "" or run.text == ""


class TestListFormattingIssue1:
    """Test list formatting fixes from issue #1."""

    def test_bullet_has_proper_indentation(self):
        """Bullet points should have proper PPTX indentation."""
        content = """## Slide

- First item
- Second item
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            # Find paragraphs with bullet chars and verify indentation
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        pPr = para._p.pPr
                        assert pPr is not None
                        buChar = pPr.find(qn('a:buChar'))
                        if buChar is not None:
                            # Check indentation is set
                            marL = pPr.get(qn('a:marL'))
                            indent = pPr.get(qn('a:indent'))
                            assert marL is not None, "marL should be set"
                            assert indent is not None, "indent should be set"

    def test_numbered_list_has_proper_indentation(self):
        """Numbered lists should have proper PPTX indentation."""
        content = """## Slide

1. First step
2. Second step
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            # Find paragraphs with auto numbering and verify indentation
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        pPr = para._p.pPr
                        assert pPr is not None
                        buAutoNum = pPr.find(qn('a:buAutoNum'))
                        if buAutoNum is not None:
                            marL = pPr.get(qn('a:marL'))
                            indent = pPr.get(qn('a:indent'))
                            assert marL is not None, "marL should be set"
                            assert indent is not None, "indent should be set"

    def test_nested_bullets_have_increasing_indentation(self):
        """Nested bullets should have increasing indentation levels."""
        content = """## Slide

- Parent item
  - Child item
    - Grandchild item
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            # Collect indentation values
            indents = []
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        pPr = para._p.pPr
                        assert pPr is not None
                        buChar = pPr.find(qn('a:buChar'))
                        if buChar is not None:
                            marL = int(pPr.get(qn('a:marL')))
                            indents.append(marL)

            assert len(indents) == 3
            # Each level should have greater indentation
            assert indents[0] < indents[1] < indents[2]

    def test_mixed_lists_work_correctly(self):
        """Mixed bullet and numbered lists should work correctly."""
        content = """## Slide

- Top-level bullet
  1. Numbered sub-item
  2. Another numbered sub-item
- Another top-level bullet
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            # Count bullets and numbers
            bullet_count = 0
            number_count = 0
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        pPr = para._p.pPr
                        assert pPr is not None
                        buChar = pPr.find(qn('a:buChar'))
                        buAutoNum = pPr.find(qn('a:buAutoNum'))
                        if buChar is not None:
                            bullet_count += 1
                        if buAutoNum is not None:
                            number_count += 1

            assert bullet_count == 2
            assert number_count == 2


class TestHangingIndentation:
    """Test PowerPoint-native hanging indentation for lists (issue #12)."""

    def test_bullet_uses_hanging_indent(self):
        """Bullet lists should use hanging indent (negative a:indent)."""
        content = """## Slide

- First item
- Second item
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        pPr = para._p.pPr
                        assert pPr is not None
                        buChar = pPr.find(qn('a:buChar'))
                        if buChar is not None:
                            indent = pPr.get(qn('a:indent'))
                            marL = pPr.get(qn('a:marL'))
                            # Hanging indent requires negative indent and positive marL
                            assert indent is not None
                            assert int(indent) < 0, f"Expected negative indent for hanging, got {indent}"
                            assert marL is not None
                            assert int(marL) > 0, f"Expected positive left margin, got {marL}"

    def test_numbered_list_uses_hanging_indent(self):
        """Numbered lists should use hanging indent (negative a:indent)."""
        content = """## Slide

1. First step
2. Second step
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        pPr = para._p.pPr
                        assert pPr is not None
                        buAutoNum = pPr.find(qn('a:buAutoNum'))
                        if buAutoNum is not None:
                            indent = pPr.get(qn('a:indent'))
                            marL = pPr.get(qn('a:marL'))
                            # Hanging indent requires negative indent and positive marL
                            assert indent is not None
                            assert int(indent) < 0, f"Expected negative indent for hanging, got {indent}"
                            assert marL is not None
                            assert int(marL) > 0, f"Expected positive left margin, got {marL}"

    def test_nested_lists_use_hanging_indent(self):
        """Nested lists should also use hanging indent (negative a:indent)."""
        content = """## Slide

- Parent bullet
  - Child bullet
1. Parent number
   1. Child number
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            indent_count = 0
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        pPr = para._p.pPr
                        assert pPr is not None
                        buChar = pPr.find(qn('a:buChar'))
                        buAutoNum = pPr.find(qn('a:buAutoNum'))
                        if buChar is not None or buAutoNum is not None:
                            indent = pPr.get(qn('a:indent'))
                            marL = pPr.get(qn('a:marL'))
                            # Hanging indent requires negative indent and positive marL
                            assert indent is not None
                            assert int(indent) < 0, f"Expected negative indent for hanging, got {indent}"
                            assert marL is not None
                            assert int(marL) > 0, f"Expected positive left margin, got {marL}"
                            indent_count += 1

            # Should have 4 list items total
            assert indent_count == 4


class TestLogoPlacement:
    """Test Multiverse Computing logo placement."""

    def test_logo_added_to_title_slide(self):
        """Title slide should have logo when logo file exists."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert("# Test Title", output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            # Check if logo path was found (depends on resources folder existing)
            if converter._logo_path:
                has_picture = any(
                    shape.shape_type == MSO_SHAPE_TYPE.PICTURE
                    for shape in slide.shapes
                )
                assert has_picture is True

    def test_logo_added_to_content_slide(self):
        """Content slide should have logo when logo file exists."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        content = """## Slide

Content here
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            if converter._logo_path:
                has_picture = any(
                    shape.shape_type == MSO_SHAPE_TYPE.PICTURE
                    for shape in slide.shapes
                )
                assert has_picture is True


class TestCosmeticChanges:
    """Test cosmetic changes from issue #4."""

    def test_subtitle_uses_dark_grey(self):
        """Subtitle should use dark grey color (#404040)."""
        content = """# Title

Subtitle text
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            found_subtitle = False
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if "Subtitle text" in run.text:
                                assert run.font.color.rgb == (0x40, 0x40, 0x40)
                                found_subtitle = True
            assert found_subtitle is True

    def test_text_has_line_spacing(self):
        """Text should have half-space line spacing (space_after = 9pt)."""
        content = """## Slide

- Item one
- Item two
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            found_spacing = False
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        # Check if space_after is set
                        if para.space_after is not None and para.space_after.pt == 9:
                            found_spacing = True
            assert found_spacing is True

    def test_text_size_is_18pt(self):
        """Body text should be 18pt (issue #4)."""
        content = """## Slide

- Body text here
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            found_text = False
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if "Body text here" in run.text:
                                assert run.font.size.pt == 18
                                found_text = True
            assert found_text is True


class TestTextScaling:
    """Test auto-scaling text to fit slide (issue #7)."""

    def test_content_frame_has_auto_size(self):
        """Content frame should have TEXT_TO_FIT_SHAPE auto size mode."""
        from pptx.enum.text import MSO_AUTO_SIZE

        content = """## Slide

- Item one
- Item two
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            # Find content shape (not the title)
            found_auto_size = False
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    tf = shape.text_frame
                    # Content frame has list items
                    all_text = "".join(r.text for p in tf.paragraphs for r in p.runs)
                    if "Item one" in all_text:
                        assert tf.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                        found_auto_size = True
            assert found_auto_size is True

    def test_long_content_fits_slide(self):
        """Long content should be contained within slide boundaries."""
        # Create content with many items that would normally overflow
        items = "\n".join([f"- Item {i} with some longer text" for i in range(20)])
        content = f"""## Slide with Many Items

{items}
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            # Should not raise any errors
            result = converter.convert(content, output_path)
            assert os.path.exists(result)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            # Verify content shape exists
            found_content = False
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    all_text = "".join(r.text for p in shape.text_frame.paragraphs for r in p.runs)
                    if "Item 0" in all_text:
                        found_content = True
                        # Verify auto_size is set for shrinking
                        from pptx.enum.text import MSO_AUTO_SIZE
                        assert shape.text_frame.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            assert found_content is True


class TestHyperlinkSupport:
    """Test hyperlink/URL support (issue #6)."""

    def test_url_with_caption_creates_hyperlink(self):
        """URL with caption should create clickable hyperlink."""
        content = """## Slide

- Visit [Google](https://google.com) for search
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            found_hyperlink = False
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if "Google" in run.text:
                                assert run.hyperlink.address == "https://google.com"
                                found_hyperlink = True
            assert found_hyperlink is True

    def test_url_without_caption_uses_url_as_text(self):
        """URL without caption should display URL as text."""
        content = """## Slide

- Check [](https://example.com)
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            found_url_text = False
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if "https://example.com" in run.text:
                                assert run.hyperlink.address == "https://example.com"
                                found_url_text = True
            assert found_url_text is True

    def test_hyperlink_uses_blue_color(self):
        """Hyperlinks should use blue color (#0066CC)."""
        content = """## Slide

- Visit [Example](https://example.com)
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            found_link = False
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if "Example" in run.text:
                                assert run.font.color.rgb == (0x00, 0x66, 0xCC)
                                found_link = True
            assert found_link is True

    def test_hyperlink_is_underlined(self):
        """Hyperlinks should be underlined."""
        content = """## Slide

- Visit [Example](https://example.com)
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            found_link = False
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if "Example" in run.text:
                                assert run.font.underline is True
                                found_link = True
            assert found_link is True

    def test_url_in_plain_text(self):
        """URLs should work in plain text paragraphs."""
        content = """## Slide

Learn more at [our website](https://multiverse.com)
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            found_link = False
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if "our website" in run.text:
                                assert run.hyperlink.address == "https://multiverse.com"
                                found_link = True
            assert found_link is True

    def test_multiple_urls_in_one_line(self):
        """Multiple URLs in the same line should all be hyperlinks."""
        content = """## Slide

- Visit [Google](https://google.com) or [Bing](https://bing.com)
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            found_google = False
            found_bing = False
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if "Google" in run.text:
                                assert run.hyperlink.address == "https://google.com"
                                found_google = True
                            if "Bing" in run.text:
                                assert run.hyperlink.address == "https://bing.com"
                                found_bing = True
            assert found_google is True
            assert found_bing is True


class TestImageSupport:
    """Test image display support (issue #5)."""

    @pytest.fixture
    def test_image_path(self, tmp_path):
        """Create a test image file."""
        from PIL import Image as PILImage

        img = PILImage.new('RGB', (400, 300), color='blue')
        img_path = tmp_path / "test_image.png"
        img.save(img_path)
        return str(img_path)

    def test_image_with_caption_renders(self, test_image_path):
        """Image with caption should render in the slide."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        content = f"""## Slide with Image

- Some text content
![Test Caption]({test_image_path})
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            # Check for picture shape (not counting logo)
            picture_count = sum(
                1 for shape in slide.shapes
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            )
            # Should have at least 1 picture (the content image)
            # May also have logo if resources/multiverse_logo.png exists
            assert picture_count >= 1

            # Check for caption text
            caption_found = False
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    all_text = "".join(r.text for p in shape.text_frame.paragraphs for r in p.runs)
                    if "Test Caption" in all_text:
                        caption_found = True
            assert caption_found is True

    def test_image_without_caption_renders(self, test_image_path):
        """Image without caption should render without caption text."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        content = f"""## Slide with Image

- Some text content
![]({test_image_path})
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            # Check for picture shape
            picture_count = sum(
                1 for shape in slide.shapes
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            )
            assert picture_count >= 1

    def test_content_width_reduced_with_image(self, test_image_path):
        """Content text box should be narrower when image is present."""
        from pptx.util import Inches

        content = f"""## Slide with Image

- Some text content
![Image]({test_image_path})
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            # Find content text box (with list items)
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    all_text = "".join(r.text for p in shape.text_frame.paragraphs for r in p.runs)
                    if "text content" in all_text:
                        # Content width should be approximately half the slide
                        # Allowing some tolerance for margins
                        expected_width = Inches(5.666)
                        actual_width = shape.width
                        # Allow 0.5 inch tolerance
                        assert abs(actual_width - expected_width) < Inches(0.5)

    def test_slide_without_image_has_full_width_content(self):
        """Slide without image should have full-width content area."""
        from pptx.util import Inches

        content = """## Slide without Image

- Some text content here
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            # Find content text box
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    all_text = "".join(r.text for p in shape.text_frame.paragraphs for r in p.runs)
                    if "text content" in all_text:
                        # Content width should be full (approximately 12.333 inches)
                        expected_width = Inches(12.333)
                        actual_width = shape.width
                        assert abs(actual_width - expected_width) < Inches(0.5)

    def test_nonexistent_image_gracefully_handled(self):
        """Nonexistent image file should not crash converter."""
        content = """## Slide with Missing Image

- Some text content
![Missing](nonexistent_image.png)
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            # Should not raise an error
            result = converter.convert(content, output_path)
            assert os.path.exists(result)

            prs = Presentation(output_path)
            assert len(prs.slides) == 1

    def test_image_caption_is_italic(self, test_image_path):
        """Image caption should be styled in italic."""
        content = f"""## Slide with Image

![My Caption]({test_image_path})
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            # Find caption and check it's italic
            found_italic_caption = False
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if "My Caption" in run.text:
                                assert run.font.italic is True
                                found_italic_caption = True
            assert found_italic_caption is True

    def test_large_image_is_scaled_down(self, tmp_path):
        """Large images should be scaled to fit the slide."""
        from PIL import Image as PILImage
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.util import Inches

        # Create a very large image
        img = PILImage.new('RGB', (2000, 1500), color='green')
        img_path = tmp_path / "large_image.png"
        img.save(img_path)

        content = f"""## Slide with Large Image

- Text content
![Large]({img_path})
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            # Find the content image (not the logo)
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    # Image should fit within right half (max ~6.2 inches)
                    max_width = Inches(6.5)
                    assert shape.width <= max_width


class TestSectionTitleStyling:
    """Test section title (H3/H4) styling (issue #8)."""

    def test_h3_section_title_is_bold_red(self):
        """H3 section titles should be bold and red (#FF0000)."""
        content = """## Slide

### Section Title

- Content
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            found_section_title = False
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if "Section Title" in run.text:
                                assert run.font.bold is True
                                assert run.font.color.rgb == (0xFF, 0x00, 0x00)
                                found_section_title = True
            assert found_section_title is True

    def test_h4_section_subtitle_is_bold_black(self):
        """H4 section subtitles should be bold and black (#111417)."""
        content = """## Slide

#### Section Subtitle

- Content
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            found_section_subtitle = False
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if "Section Subtitle" in run.text:
                                assert run.font.bold is True
                                assert run.font.color.rgb == (0x11, 0x14, 0x17)
                                found_section_subtitle = True
            assert found_section_subtitle is True

    def test_section_title_has_half_line_space_before(self):
        """Section titles should have half line space before (9pt)."""
        content = """## Slide

- Item before

### Section Title

- Item after
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            found_spacing = False
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if "Section Title" in run.text:
                                assert para.space_before is not None
                                assert para.space_before.pt == 9
                                found_spacing = True
            assert found_spacing is True

    def test_section_title_has_half_line_space_after(self):
        """Section titles should have half line space after (9pt)."""
        content = """## Slide

### Section Title

- Item after
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            found_spacing = False
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if "Section Title" in run.text:
                                assert para.space_after is not None
                                assert para.space_after.pt == 9
                                found_spacing = True
            assert found_spacing is True

    def test_h3_and_h4_different_colors(self):
        """H3 and H4 in the same slide should have different colors."""
        content = """## Slide

### H3 Title

#### H4 Title

- Content
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            h3_color = None
            h4_color = None
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if "H3 Title" in run.text:
                                h3_color = run.font.color.rgb
                            if "H4 Title" in run.text:
                                h4_color = run.font.color.rgb

            assert h3_color == (0xFF, 0x00, 0x00)  # Red
            assert h4_color == (0x11, 0x14, 0x17)  # Woodsmoke/black
            assert h3_color != h4_color

    def test_section_title_uses_body_font_size(self):
        """Section titles should use the body font size (18pt)."""
        content = """## Slide

### Section Title

- Content
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            found_section_title = False
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if "Section Title" in run.text:
                                assert run.font.size.pt == 18
                                found_section_title = True
            assert found_section_title is True

    def test_section_title_uses_body_font(self):
        """Section titles should use Open Sans font."""
        content = """## Slide

### Section Title

- Content
"""
        converter = MarkdownToPptxConverter()
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "test.pptx")
            converter.convert(content, output_path)

            prs = Presentation(output_path)
            slide = prs.slides[0]

            found_section_title = False
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if "Section Title" in run.text:
                                assert run.font.name == "Open Sans"
                                found_section_title = True
            assert found_section_title is True
