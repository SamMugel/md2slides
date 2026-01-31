"""PowerPoint converter for creating PPTX from parsed markdown."""

from __future__ import annotations

import os
from pathlib import Path
from typing import List, Optional, Union

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE

from md2slides.parser import ListItem, MarkdownParser, Slide, TextRun, ValidationError

# Bullet characters for different indentation levels
BULLET_CHARS = ['•', '–', '◦', '▪']

# Multiverse Computing Brand Colors (from style.md)
BRAND_RED = RGBColor(0xFF, 0x00, 0x00)  # #FF0000 - accent color
BRAND_CATSKILL_WHITE = RGBColor(0xF8, 0xFA, 0xFC)  # #F8FAFC - light background
BRAND_WOODSMOKE = RGBColor(0x11, 0x14, 0x17)  # #111417 - dark text
BRAND_DARK_GREY = RGBColor(0x40, 0x40, 0x40)  # #404040 - child list items

# Typography sizes (issue #4: text size to 18pt)
FONT_SIZE_TITLE = Pt(28)  # Title size
FONT_SIZE_H1 = Pt(28)  # H1 uses same as title
FONT_SIZE_H2 = Pt(28)  # H2 uses same as title
FONT_SIZE_BODY = Pt(18)  # Body text (issue #4: increase to 18pt)

# Font families (from style.md)
FONT_HEADER = "Montserrat"
FONT_BODY = "Open Sans"

# Logo settings (issue #4: make 30% smaller than previous 4.5")
LOGO_WIDTH = Inches(3.15)  # 4.5 * 0.7 = 3.15 inches (30% smaller)
LOGO_MARGIN = Inches(0.5)

# List formatting (issue #3: add 4 spaces after bullet/number)
LIST_ITEM_SPACING = "    "  # 4 spaces after bullet/number

# Line spacing (issue #4: add half space between lines)
LINE_SPACING_FACTOR = 1.5  # 1.0 = single, 1.5 = one and a half


class MarkdownToPptxConverter:
    """Convert markdown content to PowerPoint presentations."""

    def __init__(self, logo_path: Optional[str] = None) -> None:
        """Initialize the converter.

        Args:
            logo_path: Optional path to logo image file. If not provided,
                attempts to find logo in default locations.
        """
        self._prs: Presentation | None = None
        self._logo_path = self._find_logo_path(logo_path)

    def _find_logo_path(self, logo_path: Optional[str] = None) -> Optional[str]:
        """Find the logo file path.

        Args:
            logo_path: Optional explicit logo path.

        Returns:
            Path to logo file, or None if not found.
        """
        if logo_path and os.path.exists(logo_path):
            return logo_path

        # Try default locations relative to package
        package_dir = Path(__file__).parent
        default_locations = [
            package_dir / "resources" / "multiverse_logo.png",
            package_dir.parent.parent / "resources" / "multiverse_logo.png",
            Path.cwd() / "resources" / "multiverse_logo.png",
        ]

        for path in default_locations:
            if path.exists():
                return str(path)

        return None

    def convert(self, markdown_content: str, output_path: str) -> str:
        """Convert markdown content to a PPTX file.

        Args:
            markdown_content: The markdown string to convert.
            output_path: Path where the PPTX file will be saved.

        Returns:
            The absolute path to the created PPTX file.

        Raises:
            ValidationError: If markdown content or output path is invalid.
        """
        self._validate_output_path(output_path)

        parser = MarkdownParser(markdown_content)
        slides = parser.parse()

        self._prs = Presentation()
        self._prs.slide_width = Inches(13.333)
        self._prs.slide_height = Inches(7.5)

        for slide_data in slides:
            if slide_data.is_title_slide:
                self._create_title_slide(slide_data)
            else:
                self._create_content_slide(slide_data)

        # Ensure output directory exists
        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)

        self._prs.save(output_path)
        return os.path.abspath(output_path)

    def _validate_output_path(self, output_path: str) -> None:
        """Validate the output path.

        Args:
            output_path: The path to validate.

        Raises:
            ValidationError: If the path is invalid.
        """
        if not isinstance(output_path, str):
            raise ValidationError(
                f"Output path must be a string, got {type(output_path).__name__}"
            )
        if not output_path.strip():
            raise ValidationError("Output path cannot be empty")
        if not output_path.lower().endswith(".pptx"):
            raise ValidationError("Output path must have .pptx extension")

    def _add_logo_to_slide(self, slide) -> None:
        """Add the Multiverse Computing logo to a slide.

        Args:
            slide: The PowerPoint slide to add the logo to.
        """
        if not self._logo_path:
            return

        # Position logo in bottom-right corner with proper margin
        slide_width = self._prs.slide_width
        slide_height = self._prs.slide_height

        # Calculate logo position (bottom-right with margin)
        # Logo is 3.15 inches (issue #4: 30% smaller than 4.5")
        logo_left = slide_width - LOGO_WIDTH - LOGO_MARGIN
        # Estimate logo height based on aspect ratio (scales with width)
        # At 3.15" width, height is approximately 0.55" (0.78 * 0.7)
        estimated_logo_height = Inches(0.55)
        logo_top = slide_height - estimated_logo_height - LOGO_MARGIN

        slide.shapes.add_picture(
            self._logo_path,
            logo_left,
            logo_top,
            width=LOGO_WIDTH
        )

    def _set_slide_background(self, slide) -> None:
        """Set the slide background to brand color.

        Args:
            slide: The PowerPoint slide to set background for.
        """
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BRAND_CATSKILL_WHITE

    def _create_title_slide(self, slide_data: Slide) -> None:
        """Create a title slide.

        Args:
            slide_data: The slide data to render.
        """
        blank_layout = self._prs.slide_layouts[6]  # Blank layout
        slide = self._prs.slides.add_slide(blank_layout)

        # Apply brand background
        self._set_slide_background(slide)

        # Title text box - fit to text size (issue #4)
        title_left = Inches(0.5)
        title_top = Inches(2.8)  # Adjusted for better centering
        title_width = Inches(12.333)
        # Height fits the text - single line title at 28pt needs ~0.5"
        title_height = Inches(0.6)

        title_shape = slide.shapes.add_textbox(
            title_left, title_top, title_width, title_height
        )
        title_frame = title_shape.text_frame
        title_frame.word_wrap = True
        title_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_run = title_para.add_run()
        title_run.text = slide_data.title
        title_run.font.size = FONT_SIZE_H1
        title_run.font.bold = True
        title_run.font.name = FONT_HEADER
        title_run.font.color.rgb = BRAND_WOODSMOKE

        # Subtitle text box (if present) - closer to title, dark grey (issue #4)
        if slide_data.subtitle:
            # Position subtitle closer to title (0.3" gap instead of 1.7")
            subtitle_top = Inches(3.5)  # Closer to title
            subtitle_height = Inches(0.5)

            subtitle_shape = slide.shapes.add_textbox(
                title_left, subtitle_top, title_width, subtitle_height
            )
            subtitle_frame = subtitle_shape.text_frame
            subtitle_frame.word_wrap = True
            subtitle_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.alignment = PP_ALIGN.CENTER
            subtitle_run = subtitle_para.add_run()
            subtitle_run.text = slide_data.subtitle
            subtitle_run.font.size = FONT_SIZE_H2
            subtitle_run.font.name = FONT_BODY
            subtitle_run.font.color.rgb = BRAND_DARK_GREY  # Dark grey (issue #4)

        # Add logo to slide
        self._add_logo_to_slide(slide)

    def _create_content_slide(self, slide_data: Slide) -> None:
        """Create a content slide.

        Args:
            slide_data: The slide data to render.
        """
        blank_layout = self._prs.slide_layouts[6]  # Blank layout
        slide = self._prs.slides.add_slide(blank_layout)

        # Apply brand background
        self._set_slide_background(slide)

        # Title text box
        title_left = Inches(0.5)
        title_top = Inches(0.4)
        title_width = Inches(12.333)
        title_height = Inches(0.8)

        title_shape = slide.shapes.add_textbox(
            title_left, title_top, title_width, title_height
        )
        title_frame = title_shape.text_frame
        title_frame.word_wrap = True

        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.LEFT
        title_run = title_para.add_run()
        title_run.text = slide_data.title
        title_run.font.size = FONT_SIZE_H2
        title_run.font.bold = True
        title_run.font.name = FONT_HEADER
        title_run.font.color.rgb = BRAND_WOODSMOKE

        # Content text box
        content_left = Inches(0.5)
        content_top = Inches(1.4)
        content_width = Inches(12.333)
        content_height = Inches(5.6)

        content_shape = slide.shapes.add_textbox(
            content_left, content_top, content_width, content_height
        )
        content_frame = content_shape.text_frame
        content_frame.word_wrap = True
        # Auto-shrink text to fit within slide boundaries (issue #7)
        content_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        self._render_content(content_frame, slide_data.content)

        # Add logo to slide
        self._add_logo_to_slide(slide)

    def _render_content(
        self, text_frame, content: List[Union[ListItem, TextRun]]
    ) -> None:
        """Render content to a text frame.

        Args:
            text_frame: The PowerPoint text frame to render to.
            content: The list of content items to render.
        """
        first_item = True

        for item in content:
            if isinstance(item, ListItem):
                if first_item:
                    para = text_frame.paragraphs[0]
                    first_item = False
                else:
                    para = text_frame.add_paragraph()

                # Set indentation based on level
                para.level = item.level

                # Configure paragraph properties for proper list formatting
                pPr = para._p.get_or_add_pPr()

                # Remove any existing bullet settings
                for child in list(pPr):
                    tag_name = etree.QName(child.tag).localname if isinstance(child.tag, str) else ''
                    if tag_name.startswith('bu'):
                        pPr.remove(child)

                # Set indentation for proper alignment
                indent_per_level = Inches(0.5)
                left_margin = int(indent_per_level.emu * (item.level + 1))
                hanging_indent = int(Inches(0.25).emu)
                pPr.set(qn('a:marL'), str(left_margin))
                pPr.set(qn('a:indent'), str(-hanging_indent))

                if item.ordered:
                    # Numbered list using buAutoNum
                    buAutoNum = etree.SubElement(pPr, qn('a:buAutoNum'))
                    # Use different numbering styles for different levels
                    if item.level == 0:
                        buAutoNum.set('type', 'arabicPeriod')  # 1. 2. 3.
                    else:
                        buAutoNum.set('type', 'alphaLcPeriod')  # a. b. c.
                else:
                    # Bullet point using buChar
                    buChar = etree.SubElement(pPr, qn('a:buChar'))
                    bullet_char = BULLET_CHARS[min(item.level, len(BULLET_CHARS) - 1)]
                    buChar.set('char', bullet_char)

                # Add line spacing - half space between lines (issue #4)
                # Use space_after to add spacing after each paragraph
                para.space_after = Pt(9)  # Half of 18pt = 9pt extra space

                # Determine text color based on nesting level (issue #3)
                text_color = BRAND_DARK_GREY if item.level > 0 else BRAND_WOODSMOKE

                # Add 4 spaces after bullet/number (issue #3)
                spacing_run = para.add_run()
                spacing_run.text = LIST_ITEM_SPACING
                spacing_run.font.size = FONT_SIZE_BODY
                spacing_run.font.name = FONT_BODY
                spacing_run.font.color.rgb = text_color

                # Add content with formatting
                for text_run in item.content:
                    run = para.add_run()
                    run.text = text_run.text
                    run.font.size = FONT_SIZE_BODY
                    run.font.name = FONT_BODY
                    run.font.bold = text_run.bold
                    run.font.italic = text_run.italic
                    run.font.color.rgb = text_color

            elif isinstance(item, TextRun):
                if first_item:
                    para = text_frame.paragraphs[0]
                    first_item = False
                else:
                    para = text_frame.add_paragraph()

                # Add line spacing - half space between lines (issue #4)
                para.space_after = Pt(9)  # Half of 18pt = 9pt extra space

                run = para.add_run()
                run.text = item.text
                run.font.size = FONT_SIZE_BODY
                run.font.name = FONT_BODY
                run.font.bold = item.bold
                run.font.italic = item.italic
                run.font.color.rgb = BRAND_WOODSMOKE


def convert_file(input_path: str, output_path: str | None = None) -> str:
    """Convert a markdown file to PPTX.

    Args:
        input_path: Path to the markdown file.
        output_path: Optional path for the output PPTX. If not provided,
            uses the same name as input with .pptx extension.

    Returns:
        The absolute path to the created PPTX file.

    Raises:
        ValidationError: If the input file doesn't exist or is invalid.
        FileNotFoundError: If the input file doesn't exist.
    """
    if not isinstance(input_path, str):
        raise ValidationError(
            f"Input path must be a string, got {type(input_path).__name__}"
        )

    input_path = input_path.strip()
    if not input_path:
        raise ValidationError("Input path cannot be empty")

    path = Path(input_path)
    if not path.exists():
        raise FileNotFoundError(f"Input file not found: {input_path}")

    if not path.is_file():
        raise ValidationError(f"Input path is not a file: {input_path}")

    # Read markdown content
    content = path.read_text(encoding="utf-8")

    # Determine output path
    if output_path is None:
        output_path = str(path.with_suffix(".pptx"))

    converter = MarkdownToPptxConverter()
    return converter.convert(content, output_path)
